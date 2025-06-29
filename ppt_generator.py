from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

def fit_font_size(text: str, base_size: int = 20) -> Pt:
    ln = len(text)
    if ln < 300:
        return Pt(base_size)
    elif ln < 500:
        return Pt(base_size - 2)
    elif ln < 700:
        return Pt(base_size - 4)
    else:
        return Pt(base_size - 6)

def auto_linebreak(text: str, maxlen: int = 60) -> str:
    sents = re.split(r"(。|！|？|\!|\?)", text)
    grouped = ["".join(x) for x in zip(sents[0::2], sents[1::2])]
    if len(sents) % 2 != 0:
        grouped.append(sents[-1])
    result = []
    for g in grouped:
        if len(g) <= maxlen:
            result.append(g)
        else:
            forced = [g[i:i+maxlen] for i in range(0, len(g), maxlen)]
            result.extend(forced)
    return "\n".join(result)

def chunk_text(text: str, max_chars: int = 400) -> list[str]:
    paras = [p for p in text.split("\n") if p.strip()]
    pages, buf = [], ""
    for para in paras:
        if len(buf) + len(para) + 1 <= max_chars:
            buf = buf + "\n" + para if buf else para
        else:
            if buf:
                pages.append(buf)
            p = para
            while len(p) > max_chars:
                pages.append(p[:max_chars])
                p = p[max_chars:]
            buf = p
    if buf:
        pages.append(buf)
    return pages

def set_font(text_frame, font_name: str, font_size: Pt = Pt(24), bold: bool = False, align_center: bool = False, font_color: RGBColor = RGBColor(0,0,0)):
    for p in text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER if align_center else PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = font_color

def create_ppt(
    slides: list[dict],
    image_paths: list[str],
    background: str | None = None,
    title_font: str = "微软雅黑",
    body_font: str = "微软雅黑",
    color_style: str = "默认"
) -> str:
    prs = Presentation()
    w, h = prs.slide_width, prs.slide_height
    MAX_CHARS_PER_SLIDE = 400

    # 配色
    color_map = {
        "默认": RGBColor(0,0,0),
        "蓝白": RGBColor(0,0,0),
        "黑金": RGBColor(218,165,32),
        "绿色生态": RGBColor(0,128,0)
    }
    font_color = color_map.get(color_style, RGBColor(0,0,0))

    for slide in slides:
        is_image_slide = "image_path" in slide

        if is_image_slide:
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            if background:
                bg = sl.shapes.add_picture(background, 0, 0, width=w, height=h)
                spTree = sl.shapes._spTree
                sp = bg._element
                spTree.remove(sp)
                spTree.insert(2, sp)

            tb_title = sl.shapes.add_textbox(Inches(0.8), Inches(0.3), w - Inches(1.6), Inches(1))
            tf_title = tb_title.text_frame
            ani = slide.get("animation", "")
            tf_title.text = slide["title"] + (f" (推荐动画: {ani})" if ani else "")
            set_font(tf_title, title_font, Pt(32), bold=True, align_center=True, font_color=font_color)

            pic = sl.shapes.add_picture(slide["image_path"], 0, 0)
            pic.left = int((w - pic.width) / 2)
            pic.top = int((h - pic.height) / 2.5)

            tb_body = sl.shapes.add_textbox(Inches(0.8), Inches(5.2), w - Inches(1.6), Inches(2))
            tf_body = tb_body.text_frame
            desc = auto_linebreak(slide["content"].strip()[:200], 50)
            tf_body.text = desc
            set_font(tf_body, body_font, fit_font_size(desc), font_color=font_color)

            notes = sl.notes_slide.notes_text_frame
            notes.text = f"推荐动画：{ani}" if ani else ""

            extended = slide.get("extended", "").strip()
            if extended:
                pages = chunk_text(extended, MAX_CHARS_PER_SLIDE)
                for i, txt in enumerate(pages):
                    sl2 = prs.slides.add_slide(prs.slide_layouts[1])
                    if background:
                        bg2 = sl2.shapes.add_picture(background, 0, 0, width=w, height=h)
                        spTree2 = sl2.shapes._spTree
                        sp2 = bg2._element
                        spTree2.remove(sp2)
                        spTree2.insert(2, sp2)

                    tb2 = sl2.shapes.add_textbox(Inches(0.8), Inches(0.3), w - Inches(1.6), Inches(1))
                    tf2 = tb2.text_frame
                    suffix = f"（补充 {i+1}）" if len(pages) > 1 else "（补充）"
                    tf2.text = slide["title"] + suffix
                    set_font(tf2, title_font, Pt(32), bold=True, font_color=font_color)

                    ph = sl2.placeholders[1]
                    ph.text = auto_linebreak(txt, 60)
                    set_font(ph.text_frame, body_font, fit_font_size(txt), font_color=font_color)

        else:
            content = slide["content"].strip()
            pages = chunk_text(content, MAX_CHARS_PER_SLIDE)
            for i, txt in enumerate(pages):
                sl = prs.slides.add_slide(prs.slide_layouts[1])
                if background:
                    bg3 = sl.shapes.add_picture(background, 0, 0, width=w, height=h)
                    spTree3 = sl.shapes._spTree
                    sp3 = bg3._element
                    spTree3.remove(sp3)
                    spTree3.insert(2, sp3)

                tb_title = sl.shapes.add_textbox(Inches(0.8), Inches(0.3), w - Inches(1.6), Inches(1))
                tf_title = tb_title.text_frame
                tf_title.text = slide["title"] if i == 0 else f"{slide['title']}（续{ i+1 }）"
                set_font(tf_title, title_font, Pt(32), bold=True, font_color=font_color)

                ph = sl.placeholders[1]
                ph.text = auto_linebreak(txt, 60)
                set_font(ph.text_frame, body_font, fit_font_size(txt), font_color=font_color)

    out = "AutoPPT_AI.pptx"
    prs.save(out)
    return out